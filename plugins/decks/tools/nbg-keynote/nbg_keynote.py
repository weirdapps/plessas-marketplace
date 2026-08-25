#!/usr/bin/env python3
"""NBG Keynote compositor — build a cinematic dark keynote from a YAML spec.

Renders 2560x1440 slides with Pillow (full-bleed imagery, directional teal scrims,
Aptos typography, film grain) and assembles them into a 13.333x7.5in PPTX plus a
matching PDF, with per-slide speaker notes.

Keynote mode is Standard #21 in shared/presentation-style-guide.md. It is the ONLY
sanctioned exception to the light-mode NBG deck format, and only for talks delivered
live to an external or bank-wide audience. Everything else uses nbg_build.py.

Spec: shared/brand-system/keynote.md

Usage:
    python3 nbg_keynote.py deck.yaml
    python3 nbg_keynote.py deck.yaml --out ~/Downloads/202608061821_my_talk
    python3 nbg_keynote.py deck.yaml --validate      # check spec, render nothing
    python3 nbg_keynote.py deck.yaml --slides 1,4,7  # re-render a subset
"""

from __future__ import annotations

import argparse
import sys
import unicodedata
from pathlib import Path

import numpy as np
import yaml
from PIL import Image, ImageDraw, ImageFilter, ImageFont

# ============================================================ canvas constants

W, H = 2560, 1440  # render resolution (2x of 1280x720)
EMU_W, EMU_H = 12192000, 6858000  # 13.333in x 7.5in
DPI = 192  # 2560 / 13.333in

M = 155  # side gutter, both sides (0.807in)
X_RIGHT = W - M  # 2405
SAFE = 100  # nothing critical within 100px of an edge

Y_KICKER = 150
Y_TITLE = 250
Y_FOOTER_LOGO = 1300
Y_SOURCE = 1312
LOGO_H_CONTENT = 46
LOGO_H_COVER = 60
LOGO_H_BACK = 150

GRAIN_SIGMA = 4
DUO_X_RIGHT = 1360  # left edge of the right-hand stat in duo-stat

# ============================================================ palette (dark)
# Mapped from the light-mode SSOT in brand-system/colors.md. See keynote.md.

INK = (255, 255, 255)  # #FFFFFF   <- 003841  primary text, hero numbers
INK_2 = (223, 230, 230)  # #DFE6E6   <- 202020  body and support text
INK_3 = (150, 166, 168)  # #96A6A8   <- 5A5F5A  source notes, metadata
ACCENT = (0, 223, 248)  # #00DFF8   <- 007B85  kicker, emphasis, highlight
ACCENT_BAR = (0, 190, 210)  # #00BED2   <- 00ADBF  default bar fill
ACCENT_MUTE = (0, 130, 146)  # #008292   <- BEC1BE  non-highlighted bars
GROUND_TOP = (0, 22, 27)  # #00161B   <- FFFFFF  gradient top, scrim base
GROUND_BOT = (0, 56, 65)  # #003841   <- F5F8F6  gradient bottom (brand Dark Teal)
NEGATIVE = (255, 82, 99)  # #FF5263   <- AA0028  negative statistic
RULE = (120, 140, 142)  # #788C8E   <- BEC1BE  chart baseline

GLOW = np.array([0, 40, 46])  # faint cyan bloom, lower-left of the gradient

NAMED_COLORS = {
    "ink": INK,
    "white": INK,
    "ink-2": INK_2,
    "ink-3": INK_3,
    "accent": ACCENT,
    "cyan": ACCENT,
    "negative": NEGATIVE,
    "red": NEGATIVE,
}

# ============================================================ typography

# Aptos is the NBG brand font. Fall back rather than crash on a machine without it.
FONT_DIRS = [
    Path.home() / "Library/Fonts",
    Path("/Library/Fonts"),
    Path("/System/Library/Fonts/Supplemental"),
    Path("/usr/share/fonts/truetype"),
]
FONT_FILES = {
    "x": ["Aptos-ExtraBold.ttf", "Aptos-Bold.ttf", "Calibri Bold.ttf", "DejaVuSans-Bold.ttf"],
    "s": ["Aptos-SemiBold.ttf", "Aptos-Bold.ttf", "Calibri Bold.ttf", "DejaVuSans-Bold.ttf"],
    "l": ["Aptos-Light.ttf", "Aptos.ttf", "Calibri.ttf", "DejaVuSans.ttf"],
    "r": ["Aptos.ttf", "Calibri.ttf", "DejaVuSans.ttf"],
}
_font_cache: dict[tuple[str, int], ImageFont.FreeTypeFont] = {}
_font_warned: set[str] = set()


def _find_font_file(name: str) -> Path | None:
    """Locate a font file, checking one level of subdirectories as well.

    The macOS dirs hold their fonts flat, so a direct hit is tried first and
    Aptos still wins there. Linux does not: Debian and Ubuntu file fonts by
    family, so the DejaVu fallback this module already declares lives at
    /usr/share/fonts/truetype/dejavu/DejaVuSans.ttf, one level below the
    directory listed above. A flat lookup could therefore never reach it, and
    the declared fallback was unreachable on exactly the platform that needs
    it: CI had no Aptos, no Calibri and no way to find DejaVu, so it raised
    SystemExit instead of degrading.
    """
    for d in FONT_DIRS:
        p = d / name
        if p.exists():
            return p
    for d in FONT_DIRS:
        if not d.is_dir():
            continue
        for sub in sorted(d.iterdir()):
            if sub.is_dir():
                p = sub / name
                if p.exists():
                    return p
    return None


def font(kind: str, size: int) -> ImageFont.FreeTypeFont:
    key = (kind, size)
    if key in _font_cache:
        return _font_cache[key]
    for i, name in enumerate(FONT_FILES[kind]):
        p = _find_font_file(name)
        if p is not None:
            if i > 0 and kind not in _font_warned:
                _font_warned.add(kind)
                warn(f"Aptos '{kind}' weight missing, falling back to {name}")
            f = ImageFont.truetype(str(p), size)
            _font_cache[key] = f
            return f
    raise SystemExit(
        f"No usable font for weight '{kind}'. Install Aptos "
        f"(https://learn.microsoft.com/typography) or Calibri."
    )


# Greek all-caps drops the tonos but keeps the dialytika. Python's str.upper()
# keeps both, so 'η συναίνεση'.upper() renders as 'Η ΣΥΝΑΊΝΕΣΗ' — a visible typo.
TONOS = "́"
DIALYTIKA = "̈"


def greek_upper(text: str) -> str:
    out = []
    for ch in unicodedata.normalize("NFD", text.upper()):
        if ch == TONOS:
            continue
        out.append(ch)
    return unicodedata.normalize("NFC", "".join(out))


def wrap(text: str, fnt: ImageFont.FreeTypeFont, maxw: int) -> list[str]:
    lines, cur = [], ""
    for word in text.split():
        trial = (cur + " " + word).strip()
        if fnt.getlength(trial) <= maxw:
            cur = trial
        else:
            if cur:
                lines.append(cur)
            cur = word
    if cur:
        lines.append(cur)
    return lines


# ============================================================ background


def cover_fit(img: Image.Image) -> Image.Image:
    """Scale to cover the canvas, then centre-crop."""
    iw, ih = img.size
    s = max(W / iw, H / ih)
    img = img.resize((int(iw * s + 1), int(ih * s + 1)), Image.Resampling.LANCZOS)
    iw, ih = img.size
    x, y = (iw - W) // 2, (ih - H) // 2
    return img.crop((x, y, x + W, y + H))


def gradient_bg() -> Image.Image:
    top, bot = np.array(GROUND_TOP), np.array(GROUND_BOT)
    col = np.linspace(0, 1, H)[:, None]
    grad = (top[None, :] * (1 - col) + bot[None, :] * col).astype(np.uint8)
    im = np.repeat(grad[:, None, :], W, axis=1).astype(np.float32)
    yy, xx = np.mgrid[0:H, 0:W]
    d = np.sqrt(((xx - 300) / 1500) ** 2 + ((yy - 1250) / 900) ** 2)
    im += np.clip(1 - d, 0, 1)[:, :, None] * GLOW
    return Image.fromarray(np.clip(im, 0, 255).astype(np.uint8), "RGB").convert("RGBA")


SCRIM_SIDES = ("left", "right", "upperleft", "bottom", "full")


def scrim(base: Image.Image, side: str = "left", strength: float = 1.0) -> Image.Image:
    """Darken and lay a directional teal scrim so text stays legible over imagery."""
    if side not in SCRIM_SIDES:
        raise SystemExit(f"Unknown scrim '{side}'. Use one of {SCRIM_SIDES}.")
    im = np.array(base.convert("RGB")).astype(np.float32) * 0.82
    yy, xx = np.mgrid[0:H, 0:W]
    if side == "left":
        a = np.clip(1 - xx / (W * 0.72), 0, 1) ** 1.4
    elif side == "right":
        a = np.clip((xx - W * 0.28) / (W * 0.72), 0, 1) ** 1.4
    elif side == "upperleft":
        a = np.clip(1 - (xx / (W * 0.8) + yy / (H * 1.1)), 0, 1) ** 1.2
    elif side == "bottom":
        a = np.clip((yy - H * 0.35) / (H * 0.65), 0, 1) ** 1.3
    else:
        a = np.full((H, W), 0.55)
    a = np.clip(a * 0.9 * strength, 0, 1)[:, :, None]
    im = im * (1 - a) + np.array(GROUND_TOP) * a
    # bottom band so the footer always reads over imagery
    ba = np.clip((yy - H * 0.80) / (H * 0.20), 0, 1)[:, :, None] * 0.5
    im = im * (1 - ba) + np.array(GROUND_TOP) * ba
    # vignette
    d = np.sqrt(((xx - W / 2) / (W / 2)) ** 2 + ((yy - H / 2) / (H / 2)) ** 2)
    im = im * (1 - np.clip(d - 0.7, 0, 1)[:, :, None] * 0.45)
    return Image.fromarray(np.clip(im, 0, 255).astype(np.uint8), "RGB").convert("RGBA")


def grain(im: Image.Image, seed: int, amt: int = GRAIN_SIGMA) -> Image.Image:
    """Film grain. Also the thing that stops the dark gradient from banding."""
    rng = np.random.default_rng(seed)
    a = np.array(im.convert("RGB")).astype(np.float32) + rng.normal(0, amt, (H, W, 1))
    return Image.fromarray(np.clip(a, 0, 255).astype(np.uint8), "RGB").convert("RGBA")


# ============================================================ contrast guard


def _linear(c: float) -> float:
    c /= 255.0
    return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4


def luminance(rgb) -> float:
    r, g, b = (_linear(float(v)) for v in rgb[:3])
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def contrast_ratio(fg, bg_lum: float) -> float:
    fl = luminance(fg)
    hi, lo = max(fl, bg_lum), min(fl, bg_lum)
    return (hi + 0.05) / (lo + 0.05)


def _bg_luminance(canvas: Image.Image, box) -> float:
    """90th-percentile luminance under a text box — the brightest realistic patch,
    not the average, because text fails where the photograph is lightest."""
    x0, y0, x1, y1 = (max(0, int(v)) for v in box)
    x1, y1 = min(W, max(x1, x0 + 1)), min(H, max(y1, y0 + 1))
    a = np.array(canvas.convert("RGB").crop((x0, y0, x1, y1))).astype(np.float32) / 255.0
    lin = np.where(a <= 0.03928, a / 12.92, ((a + 0.055) / 1.055) ** 2.4)
    lum = 0.2126 * lin[:, :, 0] + 0.7152 * lin[:, :, 1] + 0.0722 * lin[:, :, 2]
    return float(np.percentile(lum, 90))


def ensure_contrast(canvas: Image.Image, box, fg, min_ratio: float, label: str = "") -> None:
    """Deepen a feathered patch of scrim under `box` until `fg` clears `min_ratio`.

    The patch is Gaussian-blurred at a large radius so its edge is invisible against
    photography and grain. Mutates `canvas` in place.
    """
    if contrast_ratio(fg, _bg_luminance(canvas, box)) >= min_ratio:
        return
    # A wide feather keeps the patch from reading as a rectangle, but the blur must
    # stay well inside the padding or it dilutes the alpha over the text itself.
    pad, blur = 240, 150
    x0, y0, x1, y1 = box
    for alpha in (0.20, 0.35, 0.50, 0.65, 0.80, 0.92, 1.0):
        mask = Image.new("L", (W, H), 0)
        ImageDraw.Draw(mask).rectangle(
            [x0 - pad, y0 - pad, x1 + pad, y1 + pad], fill=int(255 * alpha)
        )
        patch = Image.new("RGBA", (W, H), GROUND_TOP + (255,))
        patch.putalpha(mask.filter(ImageFilter.GaussianBlur(blur)))
        trial = Image.alpha_composite(canvas.convert("RGBA"), patch)
        if contrast_ratio(fg, _bg_luminance(trial, box)) >= min_ratio:
            canvas.paste(trial, (0, 0))
            return
    canvas.paste(trial, (0, 0))
    warn(
        f"contrast below {min_ratio}:1 even at maximum scrim{' — ' + label if label else ''}. "
        f"Use a darker photograph or move the text."
    )


def _min_ratio(size: int) -> float:
    """WCAG: 3:1 for large text (>=18pt), 4.5:1 otherwise.

    2560px spans 13.333in = 960pt, so 1px = 0.375pt and 18pt = 48px.
    """
    return 3.0 if size >= 48 else 4.5


# ============================================================ drawing helpers


def shadow_text(canvas, xy, text, fnt, fill, shadow=True, anchor=None) -> None:
    if shadow:
        lay = Image.new("RGBA", (W, H), (0, 0, 0, 0))
        ImageDraw.Draw(lay).text(
            (xy[0] + 3, xy[1] + 5), text, font=fnt, fill=(0, 0, 0, 170), anchor=anchor
        )
        canvas.alpha_composite(lay.filter(ImageFilter.GaussianBlur(7)))
    ImageDraw.Draw(canvas).text(xy, text, font=fnt, fill=fill, anchor=anchor)


def para(
    canvas,
    x: int,
    y: int,
    text,
    fnt,
    fill,
    maxw: int,
    lh: int,
    shadow=False,
    size: int = 0,
    guard=True,
    label="",
) -> int:
    """Draw a wrapped paragraph. Returns the y just past the last line."""
    if not text:
        return y
    lines = wrap(text, fnt, maxw)
    if guard:
        width = max(fnt.getlength(ln) for ln in lines)
        ensure_contrast(
            canvas,
            (x, y, x + width, y + lh * len(lines)),
            fill,
            _min_ratio(size or fnt.size),
            label or text[:40],
        )
    for ln in lines:
        shadow_text(canvas, (x, y), ln, fnt, fill, shadow)
        y += lh
    return y


def line(canvas, xy, text, fnt, fill, shadow=False, guard=True, label="") -> None:
    if not text:
        return
    if guard:
        ensure_contrast(
            canvas,
            (xy[0], xy[1], xy[0] + fnt.getlength(text), xy[1] + fnt.size * 1.3),
            fill,
            _min_ratio(fnt.size),
            label or text[:40],
        )
    shadow_text(canvas, xy, text, fnt, fill, shadow)


def kicker(canvas, text, greek=False, color=ACCENT, y=Y_KICKER) -> None:
    """The keynote's section marker: a small square plus tracked-out caps.
    Same job as the filled section pill in light mode."""
    if not text:
        return
    d = ImageDraw.Draw(canvas)
    d.rectangle([M, y + 6, M + 22, y + 28], fill=color)
    f = font("s", 30)
    caps = greek_upper(text) if greek else text.upper()
    x = float(M + 44)
    for ch in caps:
        d.text((x, y), ch, font=f, fill=color)
        x += f.getlength(ch) + 6


def place_logo(canvas, logo: Image.Image, height: int, xy) -> int:
    lg = logo.resize((int(height * logo.width / logo.height), height), Image.Resampling.LANCZOS)
    canvas.alpha_composite(lg, xy)
    return int(lg.width)


def footer(canvas, logo, source=None) -> None:
    w = place_logo(canvas, logo, LOGO_H_CONTENT, (M, Y_FOOTER_LOGO))
    if source:
        ImageDraw.Draw(canvas).text((M + w + 40, Y_SOURCE), source, font=font("r", 24), fill=INK_3)


MAX_BARS = 7
MAX_CHARTS = 2  # the exemplar keynote uses two in 16 slides; more is a business deck
BAR_W_CAP = 380  # n<=3 otherwise yields ~1000px bars: a wall, not a comparison


def draw_bars(canvas, cats, vals, highlight=None, unit="") -> None:
    if len(vals) != len(cats):
        raise SystemExit("bars: 'cats' and 'vals' must be the same length")
    if len(vals) > MAX_BARS:
        raise SystemExit(f"bars: at most {MAX_BARS} bars in keynote mode (got {len(vals)})")
    x1, base, top = X_RIGHT, 1170, 545
    x0 = float(M)
    n = len(vals)
    gap = 46 if n > 3 else 150
    bw = (x1 - x0 - gap * (n - 1)) / n
    span = x1 - x0
    if n <= 3 and bw > BAR_W_CAP:
        bw = float(BAR_W_CAP)
        x0 = x0 + (span - (bw * n + gap * (n - 1))) / 2  # centre the group
    mx = max(vals)
    d = ImageDraw.Draw(canvas)
    d.line([M, base, X_RIGHT, base], fill=RULE, width=2)
    for i, (c, v) in enumerate(zip(cats, vals, strict=True)):
        bx = x0 + i * (bw + gap)
        bh = v / mx * (base - top - 90)  # true zero baseline
        col = ACCENT if highlight == i else (ACCENT_BAR if n <= 3 else ACCENT_MUTE)
        d.rectangle([bx, base - bh, bx + bw, base], fill=col)
        d.text(
            (bx + bw / 2, base - bh - 20),
            f"{v}{unit}",
            font=font("x", 60 if n <= 3 else 46),
            fill=INK,
            anchor="mb",
        )
        d.text(
            (bx + bw / 2, base + 22),
            c,
            font=font("r", 34 if n <= 3 else 30),
            fill=INK if highlight == i else INK_2,
            anchor="ma",
        )


# ============================================================ archetypes


def _ground(spec, assets: Path):
    """Photo plus scrim, or the house gradient."""
    img = spec.get("image")
    if not img:
        return gradient_bg(), False
    p = Path(img).expanduser()
    if not p.is_absolute():
        p = assets / img
    if not p.exists():
        raise SystemExit(f"image not found: {p}")
    return scrim(
        cover_fit(Image.open(p)), spec.get("scrim", "left"), float(spec.get("scrim_strength", 1.0))
    ), True


def _color(spec, key, default):
    v = spec.get(key)
    if v is None:
        return default
    if isinstance(v, str) and v in NAMED_COLORS:
        return NAMED_COLORS[v]
    if isinstance(v, str) and len(v.lstrip("#")) == 6:
        h = v.lstrip("#")
        return tuple(int(h[i : i + 2], 16) for i in (0, 2, 4))
    raise SystemExit(f"unknown colour '{v}' for '{key}'")


def render_cover(c, s, ctx, on_image):
    shadow_text(c, (M, 300), s["title"], font("x", s.get("size", 150)), INK)
    para(
        c,
        M,
        500,
        s.get("subtitle", ""),
        font("l", 50),
        INK_2,
        1450,
        68,
        shadow=True,
        label="cover subtitle",
    )
    sp = s.get("speaker") or {}
    y = 1140
    line(c, (M, y), sp.get("name", ""), font("s", 30), ACCENT, label="speaker name")
    line(c, (M, y + 45), sp.get("role", ""), font("r", 27), INK_2, label="speaker role")
    line(c, (M, y + 92), s.get("venue", ""), font("r", 24), INK_3, label="venue")
    lg = ctx["logo"]
    w = int(LOGO_H_COVER * lg.width / lg.height)
    place_logo(c, lg, LOGO_H_COVER, (X_RIGHT - w, 1195))


def render_statement(c, s, ctx, on_image):
    kicker(c, s.get("kicker"), ctx["greek"])
    size = s.get("size", 84 if on_image else 86)
    maxw = s.get("maxw", 1500 if on_image else 2260)
    lh = int(size * 1.25)
    y = para(
        c,
        M,
        s.get("y", 470),
        s["text"],
        font("l", size),
        INK,
        maxw,
        lh,
        shadow=on_image,
        size=size,
        label="statement",
    )
    if s.get("emphasis"):
        para(
            c,
            M,
            y + int(lh * 0.35),
            s["emphasis"],
            font("l", size),
            ACCENT,
            maxw,
            lh,
            shadow=on_image,
            size=size,
            label="emphasis",
        )
    if s.get("coda"):
        line(
            c,
            (M, s.get("coda_y", 980)),
            s["coda"],
            font("l", 44),
            ACCENT,
            shadow=on_image,
            label="coda",
        )
    footer(c, ctx["logo"], s.get("source"))


def render_hero_stat(c, s, ctx, on_image):
    kicker(c, s.get("kicker"), ctx["greek"])
    # A stat can sit on the open side of a photograph rather than the gutter.
    right = s.get("align", "left") == "right"
    x = int(s.get("x", 1150 if right else M))
    # Narrow the text column over a photograph: a 'left' scrim has faded by ~x1840,
    # so a full-width line runs out of its own background and onto the image.
    if right:
        maxw_cap, maxw_sup = 1250, 1250
    elif on_image:
        maxw_cap, maxw_sup = 1700, 1750
    else:
        maxw_cap, maxw_sup = 2100, 1950
    maxw_cap = int(s.get("maxw", maxw_cap))
    maxw_sup = int(s.get("maxw", maxw_sup))
    vsize = s.get("size", 400 if on_image else 440)
    vy = s.get("y", 250 if on_image else 300)
    vcolor = _color(s, "color", INK)
    shadow_text(c, (x, vy), str(s["value"]), font("x", vsize), vcolor, shadow=on_image)
    cy = s.get("caption_y", vy + int(vsize * 1.18))
    csize = s.get("caption_size", 48 if right else (58 if on_image else 60))
    # The caption is normally accent cyan. When the number is red it goes white —
    # cyan against negative red is two loud colours fighting for the same eye.
    y = para(
        c,
        x + (6 if right else 0),
        cy,
        s.get("caption", ""),
        font("l", csize),
        _color(s, "caption_color", INK if vcolor == NEGATIVE else ACCENT),
        maxw_cap,
        int(csize * 1.25),
        shadow=on_image,
        size=csize,
        label="stat caption",
    )
    para(
        c,
        x + (6 if right else 0),
        y + 76,
        s.get("support", ""),
        font("r", 34 if right else (37 if on_image else 38)),
        INK_2,
        maxw_sup,
        47 if right else 52,
        shadow=on_image,
        label="stat support",
    )
    footer(c, ctx["logo"], s.get("source"))


def render_duo_stat(c, s, ctx, on_image):
    kicker(c, s.get("kicker"), ctx["greek"])
    para(
        c,
        M,
        Y_TITLE,
        s.get("title", ""),
        font("l", 52),
        INK,
        2250,
        64,
        shadow=on_image,
        size=52,
        label="duo title",
    )
    left, right = s["left"], s["right"]
    vsize = s.get("size", 300)
    for spec, x, default_col in ((left, M, INK), (right, DUO_X_RIGHT, ACCENT)):
        shadow_text(
            c,
            (x, 470),
            str(spec["value"]),
            font("x", vsize),
            _color(spec, "color", default_col),
            shadow=on_image,
        )
        para(
            c,
            x + 6,
            820,
            spec.get("caption", ""),
            font("r", 34),
            INK_2,
            1000,
            46,
            shadow=on_image,
            label="duo caption",
        )
    footer(c, ctx["logo"], s.get("source"))


def render_bars(c, s, ctx, on_image):
    kicker(c, s.get("kicker"), ctx["greek"])
    para(
        c,
        M,
        Y_TITLE,
        s.get("title", ""),
        font("l", 54),
        INK,
        2250,
        66,
        shadow=on_image,
        size=54,
        label="bars title",
    )
    draw_bars(c, s["cats"], s["vals"], s.get("highlight"), s.get("unit", ""))
    if s.get("axis_label"):
        line(c, (M, 1235), s["axis_label"], font("r", 28), INK_2, label="axis label")
    footer(c, ctx["logo"], s.get("source"))


def render_points(c, s, ctx, on_image):
    kicker(c, s.get("kicker"), ctx["greek"])
    para(
        c,
        M,
        Y_TITLE,
        s.get("title", ""),
        font("l", 56),
        INK,
        2250,
        68,
        shadow=on_image,
        size=56,
        label="points title",
    )
    y = s.get("y", 500)
    for p in s["points"]:
        ImageDraw.Draw(c).rectangle([M, y + 12, M + 18, y + 30], fill=ACCENT)
        y = (
            para(c, M + 42, y, p, font("r", 37), INK_2, 2050, 50, shadow=on_image, label="point")
            + 34
        )
    if s.get("takeaway"):
        line(c, (M, y + 6), s["takeaway"], font("s", 44), ACCENT, shadow=on_image, label="takeaway")
    footer(c, ctx["logo"], s.get("source"))


def render_divider(c, s, ctx, on_image):
    kicker(c, s.get("kicker"), ctx["greek"])
    x = M if s.get("align", "right") == "left" else 1290
    maxw = 2250 if x == M else 1130
    para(
        c,
        x,
        s.get("y", 560),
        s["title"],
        font("x", s.get("size", 96)),
        INK,
        maxw,
        112,
        shadow=on_image,
        size=96,
        label="divider title",
    )
    place_logo(c, ctx["logo"], LOGO_H_CONTENT, (M, Y_FOOTER_LOGO))


def render_closing(c, s, ctx, on_image):
    kicker(c, s.get("kicker"), ctx["greek"])
    size = s.get("size", 72)
    para(
        c,
        M,
        s.get("y", 330),
        s["text"],
        font("l", size),
        INK,
        2050,
        92,
        shadow=on_image,
        size=size,
        label="closing",
    )
    if s.get("coda"):
        line(
            c,
            (M, s.get("coda_y", 1130)),
            s["coda"],
            font("l", 40),
            ACCENT,
            shadow=on_image,
            label="closing coda",
        )
    place_logo(c, ctx["logo"], s.get("logo_h", 52), (M, 1290))


def render_back(c, s, ctx, on_image):
    lg = ctx["logo"]
    w = int(LOGO_H_BACK * lg.width / lg.height)
    place_logo(c, lg, LOGO_H_BACK, ((W - w) // 2, (H - LOGO_H_BACK) // 2))


RENDERERS = {
    "cover": render_cover,
    "statement": render_statement,
    "hero-stat": render_hero_stat,
    "duo-stat": render_duo_stat,
    "bars": render_bars,
    "points": render_points,
    "divider": render_divider,
    "closing": render_closing,
    "back": render_back,
}
REQUIRED_KEYS = {
    "cover": ["title"],
    "statement": ["text"],
    "hero-stat": ["value"],
    "duo-stat": ["left", "right"],
    "bars": ["cats", "vals"],
    "points": ["points"],
    "divider": ["title"],
    "closing": ["text"],
    "back": [],
}

# ============================================================ assembly


def white_wordmark(path: Path) -> Image.Image:
    """The Greek wordmark, recoloured white. Greek logo on every deck, always."""
    lg = Image.open(path).convert("RGBA")
    a = np.array(lg)
    mask = a[:, :, 3] > 8
    a[:, :, 0][mask] = a[:, :, 1][mask] = a[:, :, 2][mask] = 255
    return Image.fromarray(a)


def build_pptx(images: list[Path], notes: list[str], out: Path) -> None:
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(EMU_W), Emu(EMU_H)
    blank = prs.slide_layouts[6]
    for img, note in zip(images, notes, strict=True):
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(str(img), Emu(0), Emu(0), width=Emu(EMU_W), height=Emu(EMU_H))
        if note:
            frame = s.notes_slide.notes_text_frame
            if frame is not None:
                frame.text = note
    prs.save(str(out))


def build_pdf(images: list[Path], out: Path) -> None:
    frames = [Image.open(p).convert("RGB") for p in images]
    frames[0].save(out, save_all=True, append_images=frames[1:], resolution=DPI, quality=92)


# ============================================================ validation


def warn(msg: str) -> None:
    print(f"  warning: {msg}", file=sys.stderr)


def validate(spec: dict, assets: Path) -> list[str]:
    errors: list[str] = []
    slides = spec.get("slides") or []
    if not slides:
        errors.append("spec has no slides")
    charts = 0
    for i, s in enumerate(slides, 1):
        t = s.get("type")
        if t not in RENDERERS:
            errors.append(f"slide {i}: unknown type '{t}' (expected one of {sorted(RENDERERS)})")
            continue
        for k in REQUIRED_KEYS[t]:
            if k not in s:
                errors.append(f"slide {i} ({t}): missing required key '{k}'")
        # Standard #21: a keynote slide without a speaker note is half a slide.
        if t != "back" and not str(s.get("notes", "")).strip():
            errors.append(f"slide {i} ({t}): missing speaker notes")
        if t == "bars":
            charts += 1
            if len(s.get("vals", [])) > MAX_BARS:
                errors.append(f"slide {i}: {len(s['vals'])} bars exceeds the {MAX_BARS} limit")
        if s.get("scrim") and s["scrim"] not in SCRIM_SIDES:
            errors.append(f"slide {i}: unknown scrim '{s['scrim']}'")
        if s.get("image"):
            p = Path(s["image"]).expanduser()
            if not p.is_absolute():
                p = assets / s["image"]
            if not p.exists():
                errors.append(f"slide {i}: image not found: {p}")
    if charts > MAX_CHARTS:
        errors.append(f"{charts} chart slides; keynote mode allows at most {MAX_CHARTS}")
    return errors


# ============================================================ main


def main() -> int:
    ap = argparse.ArgumentParser(description="Build an NBG keynote from a YAML spec.")
    ap.add_argument("spec", type=Path, help="YAML deck specification")
    ap.add_argument("--out", type=Path, help="output stem (overrides meta.output)")
    ap.add_argument("--assets", type=Path, help="base directory for relative image paths")
    ap.add_argument("--validate", action="store_true", help="check the spec, render nothing")
    ap.add_argument("--slides", help="comma-separated 1-based indices to re-render")
    args = ap.parse_args()

    spec = yaml.safe_load(args.spec.read_text(encoding="utf-8"))
    meta = spec.get("meta") or {}
    assets = (args.assets or Path(meta.get("assets", args.spec.parent))).expanduser()

    errors = validate(spec, assets)
    if errors:
        print("Spec validation failed:", file=sys.stderr)
        for e in errors:
            print(f"  - {e}", file=sys.stderr)
        return 1
    if args.validate:
        print(f"OK — {len(spec['slides'])} slides, spec valid.")
        return 0

    if args.out:
        stem = args.out.expanduser()
    elif meta.get("output"):
        stem = Path(meta["output"]).expanduser()
    else:
        print("No output path: pass --out or set meta.output", file=sys.stderr)
        return 1
    stem.parent.mkdir(parents=True, exist_ok=True)
    frames_dir = stem.parent / f"{stem.name}_frames"
    frames_dir.mkdir(exist_ok=True)

    logo_path = Path(
        meta.get("logo", Path(__file__).resolve().parents[2] / "assets/nbg-logo-gr.png")
    ).expanduser()
    if not logo_path.exists():
        print(f"NBG wordmark not found: {logo_path}", file=sys.stderr)
        return 1

    ctx = {
        "logo": white_wordmark(logo_path),
        "greek": str(meta.get("language", "en")).lower().startswith("el"),
    }
    base_seed = int((spec.get("theme") or {}).get("seed", 7))
    only = {int(n) for n in args.slides.split(",")} if args.slides else None

    slides = spec["slides"]
    images, notes = [], []
    for i, s in enumerate(slides, 1):
        out_img = frames_dir / f"slide_{i:02d}.jpg"
        images.append(out_img)
        notes.append(str(s.get("notes", "")))
        if only and i not in only and out_img.exists():
            continue
        ground, on_image = _ground(s, assets)
        canvas = grain(ground, base_seed + i)  # deterministic per slide
        RENDERERS[s["type"]](canvas, s, ctx, on_image)
        # JPEG 4:4:4 — 4:2:0 smears the grain and the fine type
        canvas.convert("RGB").save(out_img, quality=92, subsampling=0)
        print(f"  slide {i:02d}  {s['type']}")

    pptx_path, pdf_path = stem.with_suffix(".pptx"), stem.with_suffix(".pdf")
    build_pptx(images, notes, pptx_path)
    build_pdf(images, pdf_path)
    print(f"\n{pptx_path}  ({len(slides)} slides)")
    print(f"{pdf_path}")
    print(f"frames: {frames_dir}  (kept, so --slides can re-render one at a time)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
